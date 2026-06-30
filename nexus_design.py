#!/home/cwatt250/AI_Agent/venv/bin/python3
"""Nexus Design Studio — full-featured local design app.

Landing page lets the user start a new prototype / slide deck / template / other,
pick wireframe or high-fidelity, and browse past designs. Each design lives in
its own folder with metadata, html, and a screenshot thumbnail. The canvas view
is a live iframe with inline editing, AI-generated tweak sliders, comment pins,
draw overlay, and exports to HTML / PDF / PPTX.
"""
from __future__ import annotations

import asyncio
import io
import json
import re
import time
import uuid
from datetime import datetime
from pathlib import Path
from threading import Lock
from typing import AsyncIterator

from fastapi import FastAPI, HTTPException, Request
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, Response, StreamingResponse
from ollama import AsyncClient
from pydantic import BaseModel

HOST = "0.0.0.0"
PORT = 11436
OLLAMA_URL = "http://localhost:11434"


def _live_model(key: str = "brain", default: str = "qwen3:8b") -> str:
    """Resolve from models.json (was hardcoded qwen3.6). Resident brain = 0 extra VRAM."""
    try:
        from pathlib import Path as _P
        return json.loads((_P.home() / "AI_Agent" / "models.json").read_text()).get(key) or default
    except Exception:
        return default


MODEL = _live_model("brain")

HOME = Path.home() / "AI_Agent"
DESIGNS_DIR = HOME / "designs"
SYSTEMS_DIR = HOME / "design-systems"
ACTIVE_FILE = SYSTEMS_DIR / "_active.json"
LEGACY_SINGLE = HOME / "design-system.json"

SSE_HEADERS = {
    "Cache-Control": "no-cache, no-transform",
    "Connection": "keep-alive",
    "X-Accel-Buffering": "no",
    "Content-Type": "text/event-stream",
}


# ---------------------------------------------------------------------------
# System prompts
# ---------------------------------------------------------------------------

_COMMON_RULES = """- Start with <!doctype html>. End with </html>. Nothing before or after.
- No markdown, no code fences, no commentary, no preamble.
- Self-contained: only load external resources from CDNs (tailwind/flowbite/reveal/google fonts).
- Always include Flowbite alongside Tailwind — Flowbite is our component library. In <head> add:
    <link href="https://cdn.jsdelivr.net/npm/flowbite@3/dist/flowbite.min.css" rel="stylesheet" />
  Before </body> add:
    <script src="https://cdn.jsdelivr.net/npm/flowbite@3/dist/flowbite.min.js"></script>
  Use Flowbite component classes wherever a component fits: buttons, cards, navbars, sidebars, modals, drawers, forms/inputs, selects, tabs, dropdowns, accordions, alerts, badges, tooltips, toggles, tables. Don't reinvent patterns Flowbite already provides.
- Use inline SVG for icons.
- Include realistic placeholder content so the design reads as a real product at a glance."""

PROMPT_PROTOTYPE_HIGH = f"""You are a senior product designer. Output ONLY a complete standalone HTML document for a polished, production-grade prototype.
{_COMMON_RULES}
- Load Tailwind via <script src="https://cdn.tailwindcss.com"></script> in <head>.
- Full visual polish: typography hierarchy, thoughtful spacing, consistent radii, good contrast.
- Use semantic HTML (header, nav, main, section, aside, footer) when it fits.
Return ONLY the HTML."""

PROMPT_PROTOTYPE_WIRE = f"""You are a senior product designer producing a LOW-FIDELITY WIREFRAME.
{_COMMON_RULES}
- Load Tailwind via <script src="https://cdn.tailwindcss.com"></script> in <head>.
- Use only grayscale: #ffffff background, #111 text, #999 secondary text, #ddd borders.
- Rectangles, dashed outlines, simple icons. Replace imagery with labeled boxes like "Image". Replace copy with realistic labels but keep everything visually muted.
- No brand colors, no gradients, no shadows beyond a single subtle border.
- Still show full structure, layout, and all sections the user asked for.
Return ONLY the HTML."""

PROMPT_SLIDES_HIGH = f"""You are a senior presentation designer. Output ONLY a complete standalone reveal.js slide deck (HTML).
{_COMMON_RULES}
- Load reveal.js via these CDN tags in <head>:
  <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5/dist/reveal.css">
  <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5/dist/theme/black.css">
- At end of <body>: <script src="https://cdn.jsdelivr.net/npm/reveal.js@5/dist/reveal.js"></script><script>Reveal.initialize({{ hash: true, controls: true, progress: true }});</script>
- Body structure: <div class="reveal"><div class="slides"><section>...</section>...</div></div>
- Produce between 6 and 10 slides: title, agenda, 3-6 content slides, summary.
- Use inline <style> in <head> to refine typography, colors, and spacing; override reveal defaults for a modern, high-end look.
Return ONLY the HTML."""

PROMPT_SLIDES_WIRE = f"""You are a senior presentation designer producing a LOW-FIDELITY WIREFRAME deck (reveal.js).
{_COMMON_RULES}
- reveal.js setup as in the high-fidelity version (CDN links, Reveal.initialize).
- Strictly grayscale, rectangles, dashed outlines, placeholder boxes for charts/images.
- 6-8 slides covering full structure.
Return ONLY the HTML."""

PROMPT_OTHER = f"""You are a senior web designer. Output ONLY a complete standalone HTML document matching the user's request.
{_COMMON_RULES}
- Load Tailwind via <script src="https://cdn.tailwindcss.com"></script> in <head> by default.
- Keep it visually polished unless the user asks otherwise.
Return ONLY the HTML."""

PROMPT_REFINE = f"""You are a senior designer refining an existing HTML document.
You will receive the current HTML and a refinement request.
{_COMMON_RULES}
- Preserve everything the user did not ask to change.
- Keep the same external resource set (Tailwind / flowbite / reveal / fonts) as the existing document.
Return ONLY the updated HTML."""

PROMPT_TWEAKS = """You analyze an HTML design and propose 4-6 quick tweaks a user could adjust with a slider or color picker.

Return ONLY a JSON array. Each item:
{
  "id": "kebab-case-id",
  "label": "Human label",
  "type": "color" | "range",
  "value": <default>,
  "min": <number>,           // range only
  "max": <number>,           // range only
  "step": <number>,          // range only
  "unit": "px"|"rem"|"",     // range only
  "css": "CSS rule(s) using {value} as placeholder"
}

Rules for "css":
- Selectors must actually match the HTML (tags, Tailwind classes, element ids).
- Use !important to win against Tailwind utilities.
- Color tweaks: e.g. `.bg-indigo-500, [class*="bg-indigo"] { background-color: {value} !important; }`.
- Range tweaks attach a unit: `body { font-size: {value}px !important; }`.
Pick tweaks with visible impact (accent color, background, radius, base font size, heading size, spacing).
Output: JSON array only, no commentary."""


# ---------------------------------------------------------------------------
# Models
# ---------------------------------------------------------------------------

class GenerateRequest(BaseModel):
    prompt: str
    name: str | None = None
    type: str = "prototype"        # prototype | slide | template | other
    fidelity: str = "high"         # wireframe | high
    design_system: str | None = None


class RefineRequest(BaseModel):
    prompt: str


class UpdateHtmlRequest(BaseModel):
    html: str


class TweaksRequest(BaseModel):
    html: str


class CommentItem(BaseModel):
    n: int
    text: str


class ApplyCommentsRequest(BaseModel):
    comments: list[CommentItem]


class DesignSystem(BaseModel):
    name: str
    primary_color: str = ""
    secondary_color: str = ""
    font: str = ""
    logo_url: str = ""
    company_name: str = ""


# ---------------------------------------------------------------------------
# App + state
# ---------------------------------------------------------------------------

app = FastAPI(title="nexus-design", version="0.5")
_client = AsyncClient(host=OLLAMA_URL)
_fs_lock = Lock()


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_FENCE_RE = re.compile(r"^\s*```(?:html)?\s*\n(.*?)\n\s*```\s*$", re.DOTALL | re.IGNORECASE)


def _strip_fences(text: str) -> str:
    m = _FENCE_RE.match(text.strip())
    if m:
        text = m.group(1)
    text = text.strip()
    lower = text.lower()
    for marker in ("<!doctype", "<html"):
        idx = lower.find(marker)
        if idx > 0:
            text = text[idx:].strip()
            lower = text.lower()
            break
    if lower.startswith("<html"):
        text = "<!doctype html>\n" + text
    return text


def _slugify(s: str, fallback: str = "design") -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", s.lower()).strip("-")
    return slug[:48] or fallback


def _new_id(name: str) -> str:
    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    return f"{ts}-{_slugify(name, 'design')}"


def _sse(event: dict) -> bytes:
    return f"data: {json.dumps(event)}\n\n".encode()


def _chunk_text(chunk) -> str:
    if isinstance(chunk, dict):
        return chunk.get("message", {}).get("content", "") or ""
    msg = getattr(chunk, "message", None)
    return getattr(msg, "content", "") or ""


def _resp_text(resp) -> str:
    if isinstance(resp, dict):
        return resp.get("message", {}).get("content", "") or ""
    msg = getattr(resp, "message", None)
    return getattr(msg, "content", "") or ""


# ---------------------------------------------------------------------------
# Design-system storage
# ---------------------------------------------------------------------------

def _ensure_systems_dir():
    SYSTEMS_DIR.mkdir(parents=True, exist_ok=True)
    if LEGACY_SINGLE.exists():
        try:
            data = json.loads(LEGACY_SINGLE.read_text())
            migrated = {
                "name": "default",
                "primary_color": data.get("primary_color", ""),
                "secondary_color": data.get("secondary_color", ""),
                "font": data.get("font", ""),
                "logo_url": "",
                "company_name": data.get("company_name", ""),
            }
            dst = SYSTEMS_DIR / "default.json"
            if not dst.exists():
                dst.write_text(json.dumps(migrated, indent=2))
            if not ACTIVE_FILE.exists():
                ACTIVE_FILE.write_text(json.dumps({"name": "default"}))
        except (OSError, json.JSONDecodeError):
            pass
        LEGACY_SINGLE.unlink(missing_ok=True)


def _list_systems() -> list[dict]:
    _ensure_systems_dir()
    out = []
    for f in sorted(SYSTEMS_DIR.glob("*.json")):
        if f.name.startswith("_"):
            continue
        try:
            out.append(json.loads(f.read_text()))
        except (OSError, json.JSONDecodeError):
            continue
    return out


def _get_active_system() -> dict | None:
    _ensure_systems_dir()
    if not ACTIVE_FILE.exists():
        return None
    try:
        name = json.loads(ACTIVE_FILE.read_text()).get("name")
    except (OSError, json.JSONDecodeError):
        return None
    if not name:
        return None
    path = SYSTEMS_DIR / f"{_slugify(name, name)}.json"
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text())
    except (OSError, json.JSONDecodeError):
        return None


def _save_system(ds: DesignSystem) -> Path:
    _ensure_systems_dir()
    path = SYSTEMS_DIR / f"{_slugify(ds.name, ds.name)}.json"
    path.write_text(json.dumps(ds.model_dump(), indent=2))
    return path


def _delete_system(name: str):
    path = SYSTEMS_DIR / f"{_slugify(name, name)}.json"
    path.unlink(missing_ok=True)
    active = _get_active_system()
    if active and active.get("name") == name:
        ACTIVE_FILE.unlink(missing_ok=True)


def _set_active_system(name: str | None):
    _ensure_systems_dir()
    if not name:
        ACTIVE_FILE.unlink(missing_ok=True)
        return
    ACTIVE_FILE.write_text(json.dumps({"name": name}))


def _brand_preamble(ds: dict | None) -> str:
    if not ds:
        return ""
    parts = []
    if ds.get("company_name"):
        parts.append(f"- Company name: {ds['company_name']}")
    if ds.get("primary_color"):
        parts.append(f"- Primary color: {ds['primary_color']}")
    if ds.get("secondary_color"):
        parts.append(f"- Secondary color: {ds['secondary_color']}")
    if ds.get("font"):
        parts.append(f"- Font family: {ds['font']} (load from Google Fonts in <head>)")
    if ds.get("logo_url"):
        parts.append(f"- Logo URL: {ds['logo_url']} (use this image where a logo is appropriate)")
    if not parts:
        return ""
    return "BRAND SYSTEM (apply consistently):\n" + "\n".join(parts) + "\n\n"


def _system_prompt_for(kind: str, fidelity: str, brand: dict | None) -> str:
    pref = _brand_preamble(brand)
    if kind == "slide":
        base = PROMPT_SLIDES_WIRE if fidelity == "wireframe" else PROMPT_SLIDES_HIGH
    elif kind == "prototype" or kind == "template":
        base = PROMPT_PROTOTYPE_WIRE if fidelity == "wireframe" else PROMPT_PROTOTYPE_HIGH
    else:
        base = PROMPT_OTHER
    return pref + base


# ---------------------------------------------------------------------------
# Design persistence
# ---------------------------------------------------------------------------

def _design_dir(id: str) -> Path:
    return DESIGNS_DIR / id


def _load_metadata(id: str) -> dict:
    meta_path = _design_dir(id) / "metadata.json"
    if not meta_path.exists():
        raise HTTPException(404, f"design not found: {id}")
    return json.loads(meta_path.read_text())


def _save_metadata(id: str, meta: dict):
    DESIGNS_DIR.mkdir(parents=True, exist_ok=True)
    (_design_dir(id)).mkdir(parents=True, exist_ok=True)
    (_design_dir(id) / "metadata.json").write_text(json.dumps(meta, indent=2))


def _load_html(id: str) -> str:
    path = _design_dir(id) / "index.html"
    if not path.exists():
        return ""
    return path.read_text(encoding="utf-8")


def _save_html(id: str, html: str):
    (_design_dir(id)).mkdir(parents=True, exist_ok=True)
    (_design_dir(id) / "index.html").write_text(html, encoding="utf-8")


def _list_designs() -> list[dict]:
    if not DESIGNS_DIR.exists():
        return []
    out = []
    for d in sorted(DESIGNS_DIR.iterdir(), key=lambda p: p.stat().st_mtime, reverse=True):
        if not d.is_dir():
            continue
        meta_path = d / "metadata.json"
        if not meta_path.exists():
            continue
        try:
            m = json.loads(meta_path.read_text())
        except (OSError, json.JSONDecodeError):
            continue
        m["has_thumbnail"] = (d / "thumbnail.png").exists()
        out.append(m)
    return out


# ---------------------------------------------------------------------------
# Streaming generation
# ---------------------------------------------------------------------------

async def _stream_ollama(system: str, user: str):
    """Async generator over (stage, payload) events."""
    yield ("thinking", {})
    buf: list[str] = []
    tokens = 0
    try:
        stream = await _client.chat(
            model=MODEL,
            messages=[
                {"role": "system", "content": system},
                {"role": "user", "content": user},
            ],
            stream=True,
            think=False,
            options={"temperature": 0.7, "num_ctx": 32768, "num_predict": 16384},
        )
        async for chunk in stream:
            piece = _chunk_text(chunk)
            if not piece:
                continue
            buf.append(piece)
            tokens += 1
            yield ("writing", {"tokens": tokens, "chunk": piece})
    except Exception as exc:
        yield ("error", {"message": f"{type(exc).__name__}: {exc}"})
        return
    raw = "".join(buf)
    html = _strip_fences(raw)
    yield ("raw-done", {"tokens": tokens, "html": html})


# ---------------------------------------------------------------------------
# Playwright: thumbnail + PDF + slide screenshots
# ---------------------------------------------------------------------------

async def _screenshot_html(html_path: Path, out_path: Path, width: int = 1280, height: int = 800):
    from playwright.async_api import async_playwright
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        try:
            ctx = await browser.new_context(viewport={"width": width, "height": height}, device_scale_factor=1)
            page = await ctx.new_page()
            await page.goto(html_path.as_uri(), wait_until="load", timeout=15000)
            await page.wait_for_timeout(600)
            await page.screenshot(path=str(out_path), type="png", full_page=False)
        finally:
            await browser.close()


async def _render_pdf(html_path: Path, out_path: Path):
    from playwright.async_api import async_playwright
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        try:
            page = await browser.new_page()
            await page.goto(html_path.as_uri(), wait_until="load", timeout=20000)
            await page.wait_for_timeout(600)
            await page.emulate_media(media="print")
            await page.pdf(path=str(out_path), format="Letter", print_background=True)
        finally:
            await browser.close()


async def _render_slide_pngs(html_path: Path, out_dir: Path, max_slides: int = 60) -> list[Path]:
    """Navigate through a reveal.js deck and screenshot each slide at 1920x1080."""
    from playwright.async_api import async_playwright
    out_dir.mkdir(parents=True, exist_ok=True)
    pngs: list[Path] = []
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        try:
            ctx = await browser.new_context(viewport={"width": 1920, "height": 1080})
            page = await ctx.new_page()
            await page.goto(html_path.as_uri(), wait_until="load", timeout=20000)
            await page.wait_for_timeout(800)
            total = await page.evaluate(
                "() => (window.Reveal && typeof Reveal.getTotalSlides === 'function') ? Reveal.getTotalSlides() : document.querySelectorAll('.reveal .slides > section').length"
            )
            total = max(1, min(int(total or 1), max_slides))
            for i in range(total):
                await page.evaluate(f"() => {{ if (window.Reveal) Reveal.slide({i}, 0); }}")
                await page.wait_for_timeout(400)
                img = out_dir / f"slide-{i+1:03d}.png"
                await page.screenshot(path=str(img), type="png", full_page=False)
                pngs.append(img)
        finally:
            await browser.close()
    return pngs


async def _make_thumbnail(id: str):
    try:
        html_path = _design_dir(id) / "index.html"
        out = _design_dir(id) / "thumbnail.png"
        if html_path.exists():
            await _screenshot_html(html_path, out, width=1280, height=800)
    except Exception:
        pass


def _build_pptx(slide_pngs: list[Path], out_path: Path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for img in slide_pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(out_path))


# ---------------------------------------------------------------------------
# API: designs
# ---------------------------------------------------------------------------

@app.get("/api/designs")
async def api_list_designs():
    return {"designs": _list_designs()}


@app.get("/api/designs/{id}")
async def api_get_design(id: str):
    meta = _load_metadata(id)
    meta["html"] = _load_html(id)
    return meta


@app.get("/api/designs/{id}/thumbnail.png")
async def api_thumbnail(id: str):
    path = _design_dir(id) / "thumbnail.png"
    if not path.exists():
        raise HTTPException(404, "no thumbnail yet")
    return FileResponse(str(path), media_type="image/png")


@app.delete("/api/designs/{id}")
async def api_delete_design(id: str):
    import shutil
    d = _design_dir(id)
    if not d.exists():
        raise HTTPException(404, "not found")
    shutil.rmtree(d, ignore_errors=True)
    return {"ok": True}


@app.post("/api/designs/{id}/html")
async def api_update_html(id: str, req: UpdateHtmlRequest):
    meta = _load_metadata(id)
    _save_html(id, req.html)
    meta["updated_at"] = int(time.time())
    _save_metadata(id, meta)
    asyncio.create_task(_make_thumbnail(id))
    return {"ok": True}


# ---------------------------------------------------------------------------
# API: generate (new design) and refine (existing design)
# ---------------------------------------------------------------------------

async def _sse_generate_stream(
    id: str, meta: dict, system: str, user: str, kind_history_label: str
) -> AsyncIterator[bytes]:
    yield _sse({"stage": "start", "id": id, "meta": meta})
    last_tokens = 0
    async for stage, payload in _stream_ollama(system, user):
        if stage == "thinking":
            yield _sse({"stage": "thinking"})
        elif stage == "writing":
            last_tokens = payload["tokens"]
            yield _sse({"stage": "writing", "tokens": last_tokens, "chunk": payload["chunk"]})
        elif stage == "error":
            yield _sse({"stage": "error", "message": payload["message"]})
            return
        elif stage == "raw-done":
            html = payload["html"]
            _save_html(id, html)
            meta = _load_metadata(id)
            hist = meta.setdefault("history", [])
            hist.append({"ts": int(time.time()), "kind": kind_history_label, "prompt": user_preview(user)})
            meta["updated_at"] = int(time.time())
            _save_metadata(id, meta)
            asyncio.create_task(_make_thumbnail(id))
            yield _sse({"stage": "done", "tokens": payload["tokens"], "html": html, "id": id})


def user_preview(u: str) -> str:
    u = u.strip()
    return u[:300]


@app.post("/api/generate")
async def api_generate(req: GenerateRequest):
    prompt = req.prompt.strip()
    if not prompt:
        raise HTTPException(400, "empty prompt")
    kind = req.type if req.type in {"prototype", "slide", "template", "other"} else "prototype"
    fidelity = "wireframe" if req.fidelity == "wireframe" else "high"
    brand = None
    if req.design_system:
        path = SYSTEMS_DIR / f"{_slugify(req.design_system, req.design_system)}.json"
        if path.exists():
            try:
                brand = json.loads(path.read_text())
            except (OSError, json.JSONDecodeError):
                brand = None
    if not brand:
        brand = _get_active_system()
    system = _system_prompt_for(kind, fidelity, brand)
    name = (req.name or "").strip() or prompt[:40]
    id = _new_id(name)
    meta = {
        "id": id,
        "name": name,
        "type": kind,
        "fidelity": fidelity,
        "prompt": prompt,
        "design_system": (brand or {}).get("name") if brand else None,
        "created_at": int(time.time()),
        "updated_at": int(time.time()),
        "history": [{"ts": int(time.time()), "kind": "create", "prompt": prompt}],
    }
    _save_metadata(id, meta)
    return StreamingResponse(
        _sse_generate_stream(id, meta, system, prompt, "create"),
        media_type="text/event-stream",
        headers=SSE_HEADERS,
    )


@app.post("/api/designs/{id}/refine")
async def api_refine(id: str, req: RefineRequest):
    prompt = req.prompt.strip()
    if not prompt:
        raise HTTPException(400, "empty prompt")
    meta = _load_metadata(id)
    current = _load_html(id)
    if not current:
        raise HTTPException(400, "design has no HTML yet")
    user_msg = f"CURRENT HTML:\n---\n{current}\n---\n\nREFINEMENT REQUEST:\n{prompt}"
    brand_name = meta.get("design_system")
    brand = None
    if brand_name:
        path = SYSTEMS_DIR / f"{_slugify(brand_name, brand_name)}.json"
        if path.exists():
            try:
                brand = json.loads(path.read_text())
            except (OSError, json.JSONDecodeError):
                brand = None
    system = _brand_preamble(brand) + PROMPT_REFINE
    return StreamingResponse(
        _sse_generate_stream(id, meta, system, user_msg, "refine"),
        media_type="text/event-stream",
        headers=SSE_HEADERS,
    )


@app.post("/api/designs/{id}/apply-comments")
async def api_apply_comments(id: str, req: ApplyCommentsRequest):
    meta = _load_metadata(id)
    current = _load_html(id)
    if not current:
        raise HTTPException(400, "design has no HTML yet")
    notes = [c for c in req.comments if c.text.strip()]
    if not notes:
        raise HTTPException(400, "no comment text")
    bullet = "\n".join(f"{c.n}. {c.text.strip()}" for c in notes)
    user_msg = (
        f"CURRENT HTML:\n---\n{current}\n---\n\n"
        f"APPLY THE FOLLOWING NUMBERED FEEDBACK (each pinned to the design):\n{bullet}\n\n"
        "Return the full updated HTML."
    )
    system = PROMPT_REFINE
    return StreamingResponse(
        _sse_generate_stream(id, meta, system, user_msg, "comments"),
        media_type="text/event-stream",
        headers=SSE_HEADERS,
    )


# ---------------------------------------------------------------------------
# API: tweaks
# ---------------------------------------------------------------------------

@app.post("/api/tweaks")
async def api_tweaks(req: TweaksRequest):
    html = req.html.strip()
    if not html:
        raise HTTPException(400, "no html")
    user_msg = f"HTML:\n---\n{html[:12000]}\n---\n\nReturn 4-6 tweaks as JSON."
    try:
        resp = await _client.chat(
            model=MODEL,
            messages=[
                {"role": "system", "content": PROMPT_TWEAKS},
                {"role": "user", "content": user_msg},
            ],
            stream=False,
            think=False,
            format="json",
            options={"temperature": 0.2, "num_ctx": 16384, "num_predict": 2048},
        )
    except Exception as exc:
        raise HTTPException(502, f"ollama: {type(exc).__name__}: {exc}")
    raw = _resp_text(resp).strip()
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        m = re.search(r"\[[\s\S]*\]", raw)
        if not m:
            raise HTTPException(502, f"tweaks: non-json response: {raw[:200]}")
        data = json.loads(m.group(0))
    if isinstance(data, dict):
        for key in ("tweaks", "items", "data", "list"):
            if key in data and isinstance(data[key], list):
                data = data[key]
                break
    if not isinstance(data, list):
        raise HTTPException(502, f"tweaks: expected list, got {type(data).__name__}")
    return {"tweaks": data[:6]}


# ---------------------------------------------------------------------------
# API: design systems
# ---------------------------------------------------------------------------

@app.get("/api/design-systems")
async def api_list_systems():
    active = _get_active_system()
    return {
        "systems": _list_systems(),
        "active": active.get("name") if active else None,
    }


@app.post("/api/design-systems")
async def api_save_system(ds: DesignSystem):
    if not ds.name.strip():
        raise HTTPException(400, "name required")
    _save_system(ds)
    return {"ok": True, "name": ds.name}


@app.delete("/api/design-systems/{name}")
async def api_delete_system(name: str):
    _delete_system(name)
    return {"ok": True}


@app.post("/api/design-systems/{name}/activate")
async def api_activate_system(name: str):
    path = SYSTEMS_DIR / f"{_slugify(name, name)}.json"
    if not path.exists():
        raise HTTPException(404, "not found")
    _set_active_system(name)
    return {"ok": True}


@app.delete("/api/design-systems/_active")
async def api_clear_active():
    _set_active_system(None)
    return {"ok": True}


# ---------------------------------------------------------------------------
# API: exports
# ---------------------------------------------------------------------------

@app.get("/api/designs/{id}/export/html")
async def export_html(id: str):
    html = _load_html(id)
    if not html:
        raise HTTPException(404, "no html")
    meta = _load_metadata(id)
    fname = f"{_slugify(meta.get('name','design'), 'design')}-{id}.html"
    return Response(
        content=html, media_type="text/html",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'},
    )


@app.get("/api/designs/{id}/export/pdf")
async def export_pdf(id: str):
    meta = _load_metadata(id)
    html_path = _design_dir(id) / "index.html"
    if not html_path.exists():
        raise HTTPException(404, "no html")
    out = _design_dir(id) / "export.pdf"
    try:
        await _render_pdf(html_path, out)
    except Exception as exc:
        raise HTTPException(500, f"pdf render failed: {type(exc).__name__}: {exc}")
    fname = f"{_slugify(meta.get('name','design'), 'design')}-{id}.pdf"
    return FileResponse(str(out), media_type="application/pdf", filename=fname)


@app.get("/api/designs/{id}/export/pptx")
async def export_pptx(id: str):
    meta = _load_metadata(id)
    if meta.get("type") != "slide":
        raise HTTPException(400, "pptx export is only available for slide decks")
    html_path = _design_dir(id) / "index.html"
    if not html_path.exists():
        raise HTTPException(404, "no html")
    tmp_dir = _design_dir(id) / "_slides"
    try:
        pngs = await _render_slide_pngs(html_path, tmp_dir)
    except Exception as exc:
        raise HTTPException(500, f"slide capture failed: {type(exc).__name__}: {exc}")
    if not pngs:
        raise HTTPException(500, "no slides captured")
    out = _design_dir(id) / "export.pptx"
    try:
        _build_pptx(pngs, out)
    except Exception as exc:
        raise HTTPException(500, f"pptx build failed: {type(exc).__name__}: {exc}")
    fname = f"{_slugify(meta.get('name','slides'), 'slides')}-{id}.pptx"
    return FileResponse(
        str(out),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=fname,
    )


# ---------------------------------------------------------------------------
# Shared shell head — Tailwind Play CDN + Honda-red primary + Inter + Flowbite
# ---------------------------------------------------------------------------

SHELL_HEAD = r"""
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap" rel="stylesheet">
<script src="https://cdn.tailwindcss.com"></script>
<script>
  tailwind.config = {
    theme: {
      extend: {
        colors: {
          primary: {
            50:  '#fff1f1',
            100: '#ffe2e2',
            200: '#ffc8c8',
            300: '#ff9c9c',
            400: '#ff6060',
            500: '#ff2828',
            600: '#cc0000',
            700: '#ab0505',
            800: '#8d0909',
            900: '#760d0d',
          }
        },
        fontFamily: { sans: ['Inter', 'ui-sans-serif', 'system-ui', 'sans-serif'] },
      }
    }
  }
</script>
<link href="https://cdn.jsdelivr.net/npm/flowbite@3/dist/flowbite.min.css" rel="stylesheet" />
<style>
  html, body { font-family: Inter, ui-sans-serif, system-ui, sans-serif; }
  /* Flowbite form inputs default to dark mode; hard-pin everything to light. */
  body { color-scheme: light; }
</style>
"""


# ---------------------------------------------------------------------------
# Landing page
# ---------------------------------------------------------------------------

LANDING_HTML = r"""<!doctype html>
<html lang="en" class="h-full">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>Nexus Design</title>
__SHELL_HEAD__
</head>
<body class="h-full bg-gray-50 text-gray-900 antialiased">
<div class="flex h-full">
  <!-- Sidebar -->
  <aside class="w-64 bg-white border-r border-gray-200 flex flex-col shrink-0">
    <div class="h-14 px-5 flex items-center gap-2 border-b border-gray-200">
      <span class="inline-block w-2.5 h-2.5 rounded-full bg-primary-600 shadow-[0_0_12px_rgba(204,0,0,0.55)]"></span>
      <span class="font-semibold text-gray-900 tracking-tight">Nexus Design</span>
    </div>
    <div class="p-3">
      <button id="focusPrompt" class="w-full text-white bg-primary-600 hover:bg-primary-700 focus:ring-4 focus:outline-none focus:ring-primary-200 font-medium rounded-lg text-sm px-4 py-2.5 transition inline-flex items-center justify-center gap-2">
        <svg class="w-4 h-4" fill="none" stroke="currentColor" stroke-width="2.2" viewBox="0 0 24 24"><path d="M12 5v14M5 12h14"/></svg>
        New Project
      </button>
    </div>
    <ul class="px-3 pb-2 space-y-1">
      <li><button class="tab w-full flex items-center gap-3 px-3 py-2 text-sm rounded-lg font-medium" data-type="prototype">
        <svg class="w-4 h-4 shrink-0" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><rect x="3" y="3" width="18" height="18" rx="2"/><path d="M3 9h18M9 21V9"/></svg>
        <span>Prototype</span>
      </button></li>
      <li><button class="tab w-full flex items-center gap-3 px-3 py-2 text-sm rounded-lg font-medium" data-type="slide">
        <svg class="w-4 h-4 shrink-0" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><rect x="2" y="4" width="20" height="14" rx="2"/><path d="M8 22h8M12 18v4"/></svg>
        <span>Slide deck</span>
      </button></li>
      <li><button class="tab w-full flex items-center gap-3 px-3 py-2 text-sm rounded-lg font-medium" data-type="template">
        <svg class="w-4 h-4 shrink-0" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><path d="M4 4h7v7H4zM13 4h7v4h-7zM13 11h7v9h-7zM4 14h7v6H4z"/></svg>
        <span>From template</span>
      </button></li>
      <li><button class="tab w-full flex items-center gap-3 px-3 py-2 text-sm rounded-lg font-medium" data-type="other">
        <svg class="w-4 h-4 shrink-0" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><circle cx="12" cy="12" r="9"/><path d="M9 10.5c0-1.5 1.3-2.5 3-2.5s3 1 3 2.5c0 1.2-.9 1.8-1.8 2.2-.7.3-1.2.7-1.2 1.3M12 17h.01"/></svg>
        <span>Other</span>
      </button></li>
    </ul>
    <div class="mt-auto border-t border-gray-200 p-4 flex-1 overflow-y-auto min-h-0">
      <div class="flex items-center justify-between mb-3">
        <h4 class="text-xs font-semibold text-gray-500 uppercase tracking-wider">Design Systems</h4>
        <button id="dsNew" class="text-xs text-primary-600 hover:text-primary-700 font-medium">+ new</button>
      </div>
      <div id="dsList" class="space-y-1"></div>
    </div>
  </aside>

  <!-- Main -->
  <main class="flex-1 overflow-y-auto">
    <div class="max-w-5xl mx-auto px-10 pt-12 pb-16">
      <h1 id="greet" class="text-3xl font-bold text-gray-900 tracking-tight">Good to see you.</h1>
      <p class="text-gray-500 mt-1 mb-8">Describe what you want to design. Nexus generates a polished HTML document you can edit live.</p>

      <!-- Generate form card -->
      <div class="bg-white border border-gray-200 rounded-xl p-5 shadow-sm">
        <div class="flex items-center gap-3 flex-wrap">
          <div id="fidelity" class="inline-flex rounded-lg border border-gray-200 overflow-hidden shadow-sm" role="group">
            <button data-v="high" class="fid px-4 py-2 text-sm font-medium">High fidelity</button>
            <button data-v="wireframe" class="fid px-4 py-2 text-sm font-medium border-l border-gray-200">Wireframe</button>
          </div>
          <div class="flex-1 min-w-[220px]">
            <input type="text" id="name" placeholder="Project name (optional)" class="bg-white border border-gray-200 text-sm text-gray-900 rounded-lg focus:ring-primary-500 focus:border-primary-500 block w-full px-3 py-2.5">
          </div>
        </div>
        <textarea id="prompt" rows="4" placeholder="A clean SaaS dashboard with a sidebar, KPI cards, and a data table..." class="mt-3 block w-full text-sm text-gray-900 bg-white rounded-lg border border-gray-200 focus:ring-primary-500 focus:border-primary-500 p-3 resize-y"></textarea>
        <div class="flex items-center justify-between mt-3 flex-wrap gap-2">
          <p class="text-xs text-gray-500">Active design system: <span id="activeDs" class="font-semibold text-gray-900">(none)</span></p>
          <button id="genBtn" class="text-white bg-primary-600 hover:bg-primary-700 focus:ring-4 focus:outline-none focus:ring-primary-200 font-medium rounded-lg text-sm px-5 py-2.5 inline-flex items-center gap-2 transition">
            Generate
            <svg class="w-4 h-4" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><path d="M5 12h14M13 5l7 7-7 7"/></svg>
          </button>
        </div>
      </div>

      <!-- Recent -->
      <div class="mt-12">
        <h3 class="text-xs font-semibold text-gray-500 uppercase tracking-wider mb-4">Recent designs</h3>
        <div id="recent" class="grid gap-4" style="grid-template-columns:repeat(auto-fill,minmax(260px,1fr));"></div>
      </div>
    </div>
  </main>
</div>

<!-- Generation overlay -->
<div id="overlay" class="fixed inset-0 bg-gray-900/80 hidden items-center justify-center z-50 backdrop-blur-sm">
  <div class="text-center space-y-4">
    <div class="w-14 h-14 border-4 border-gray-700 border-t-primary-500 rounded-full animate-spin mx-auto"></div>
    <div id="stageTxt" class="text-white text-sm font-medium">Thinking...</div>
    <div id="countTxt" class="text-gray-400 text-xs font-mono">0 tokens</div>
  </div>
</div>

<!-- Design system modal -->
<div id="dsModal" class="fixed inset-0 bg-black/50 hidden items-center justify-center z-50 p-4">
  <div class="bg-white rounded-xl w-[460px] max-w-full shadow-2xl">
    <div class="flex items-center justify-between px-5 pt-5 pb-3 border-b border-gray-200">
      <h3 id="dsModalTitle" class="text-lg font-semibold text-gray-900">New Design System</h3>
      <button id="dsX" class="text-gray-400 hover:text-gray-600 text-xl leading-none">×</button>
    </div>
    <div class="px-5 py-4 space-y-3">
      <div>
        <label class="block text-xs font-semibold text-gray-600 mb-1 uppercase tracking-wider">Name</label>
        <input type="text" id="ds_name" placeholder="Default" class="bg-white border border-gray-200 text-sm rounded-lg focus:ring-primary-500 focus:border-primary-500 block w-full px-3 py-2.5 text-gray-900">
      </div>
      <div class="grid grid-cols-2 gap-3">
        <div>
          <label class="block text-xs font-semibold text-gray-600 mb-1 uppercase tracking-wider">Primary</label>
          <div class="flex gap-2">
            <input type="color" id="ds_primary" value="#cc0000" class="w-10 h-[42px] border border-gray-200 rounded-lg cursor-pointer p-0.5 bg-white">
            <input type="text" id="ds_primary_t" placeholder="#CC0000" class="bg-white border border-gray-200 text-sm rounded-lg focus:ring-primary-500 focus:border-primary-500 block flex-1 min-w-0 px-3 py-2.5 text-gray-900">
          </div>
        </div>
        <div>
          <label class="block text-xs font-semibold text-gray-600 mb-1 uppercase tracking-wider">Secondary</label>
          <div class="flex gap-2">
            <input type="color" id="ds_secondary" value="#111827" class="w-10 h-[42px] border border-gray-200 rounded-lg cursor-pointer p-0.5 bg-white">
            <input type="text" id="ds_secondary_t" placeholder="#111827" class="bg-white border border-gray-200 text-sm rounded-lg focus:ring-primary-500 focus:border-primary-500 block flex-1 min-w-0 px-3 py-2.5 text-gray-900">
          </div>
        </div>
      </div>
      <div>
        <label class="block text-xs font-semibold text-gray-600 mb-1 uppercase tracking-wider">Font</label>
        <input type="text" id="ds_font" placeholder="Inter, Manrope, ..." class="bg-white border border-gray-200 text-sm rounded-lg focus:ring-primary-500 focus:border-primary-500 block w-full px-3 py-2.5 text-gray-900">
      </div>
      <div>
        <label class="block text-xs font-semibold text-gray-600 mb-1 uppercase tracking-wider">Company</label>
        <input type="text" id="ds_company" placeholder="Company name" class="bg-white border border-gray-200 text-sm rounded-lg focus:ring-primary-500 focus:border-primary-500 block w-full px-3 py-2.5 text-gray-900">
      </div>
      <div>
        <label class="block text-xs font-semibold text-gray-600 mb-1 uppercase tracking-wider">Logo URL</label>
        <input type="text" id="ds_logo" placeholder="https://..." class="bg-white border border-gray-200 text-sm rounded-lg focus:ring-primary-500 focus:border-primary-500 block w-full px-3 py-2.5 text-gray-900">
      </div>
    </div>
    <div class="flex items-center gap-2 px-5 py-4 border-t border-gray-200 bg-gray-50 rounded-b-xl">
      <button id="dsDelete" class="text-red-700 border border-red-200 hover:bg-red-50 font-medium rounded-lg text-sm px-4 py-2 mr-auto">Delete</button>
      <button id="dsCancel" class="text-gray-700 bg-white border border-gray-200 hover:bg-gray-100 font-medium rounded-lg text-sm px-4 py-2">Cancel</button>
      <button id="dsSave" class="text-white bg-primary-600 hover:bg-primary-700 font-medium rounded-lg text-sm px-4 py-2">Save</button>
    </div>
  </div>
</div>

<!-- Toast -->
<div id="toast" class="fixed bottom-5 right-5 bg-gray-900 text-white text-sm px-4 py-2.5 rounded-lg shadow-xl opacity-0 translate-y-2 transition-all duration-200 pointer-events-none z-[60] font-medium"></div>

<script>
const $ = (id) => document.getElementById(id);
const nf = new Intl.NumberFormat();
let selectedType = "prototype";
let selectedFidelity = "high";

function toast(msg, kind){
  const t = $("toast"); t.textContent = msg;
  t.classList.remove("bg-gray-900","bg-red-600","bg-green-600","opacity-0","translate-y-2","pointer-events-none");
  t.classList.add(kind==="err" ? "bg-red-600" : kind==="ok" ? "bg-green-600" : "bg-gray-900", "opacity-100","translate-y-0");
  clearTimeout(toast._t);
  toast._t = setTimeout(() => {
    t.classList.remove("opacity-100","translate-y-0");
    t.classList.add("opacity-0","translate-y-2","pointer-events-none");
  }, 2400);
}
function fmtDate(ts){ const d = new Date(ts*1000); return d.toLocaleDateString(undefined, { month:"short", day:"numeric" }); }
function escHtml(s){ return String(s).replace(/[&<>\"']/g, c => ({"&":"&amp;","<":"&lt;",">":"&gt;","\"":"&quot;","'":"&#39;"}[c])); }
function showModal(el){ el.classList.remove("hidden"); el.classList.add("flex"); }
function hideModal(el){ el.classList.add("hidden"); el.classList.remove("flex"); }

/* tabs (type) */
function paintTabs(){
  document.querySelectorAll(".tab").forEach(t => {
    const active = t.dataset.type === selectedType;
    t.classList.remove("bg-primary-50","text-primary-700","text-gray-700","hover:bg-gray-100");
    t.classList.add(active ? "bg-primary-50" : "hover:bg-gray-100", active ? "text-primary-700" : "text-gray-700");
  });
}
document.querySelectorAll(".tab").forEach(el => el.addEventListener("click", () => { selectedType = el.dataset.type; paintTabs(); }));

/* fidelity */
function paintFid(){
  document.querySelectorAll(".fid").forEach(b => {
    const on = b.dataset.v === selectedFidelity;
    b.classList.remove("bg-primary-600","text-white","bg-white","text-gray-700","hover:bg-gray-50");
    b.classList.add(on ? "bg-primary-600" : "bg-white", on ? "text-white" : "text-gray-700", on ? "" : "hover:bg-gray-50");
  });
}
document.querySelectorAll(".fid").forEach(b => b.addEventListener("click", () => { selectedFidelity = b.dataset.v; paintFid(); }));

$("focusPrompt").addEventListener("click", () => $("prompt").focus());
$("prompt").addEventListener("keydown", (e) => { if ((e.metaKey || e.ctrlKey) && e.key === "Enter"){ e.preventDefault(); $("genBtn").click(); }});

async function refreshRecent(){
  const r = await fetch("/api/designs"); const j = await r.json();
  const g = $("recent");
  if (!j.designs.length){
    g.innerHTML = '<div class="col-span-full text-sm text-gray-500 border border-dashed border-gray-300 rounded-xl p-6 text-center">No designs yet — describe something above to get started.</div>';
    return;
  }
  g.innerHTML = j.designs.map(d => {
    const thumb = d.has_thumbnail ? `/api/designs/${d.id}/thumbnail.png?t=${d.updated_at||0}` : "";
    const style = thumb ? `background-image:url('${thumb}'); background-size:cover; background-position:center;` : "background:#f3f4f6;";
    const fid = (d.fidelity === "wireframe") ? `<span class="inline-block text-[10px] uppercase tracking-wider bg-gray-100 text-gray-500 px-1.5 py-0.5 rounded mr-1">wire</span>` : "";
    return `
      <button type="button" data-id="${d.id}" class="card text-left bg-white border border-gray-200 rounded-xl overflow-hidden cursor-pointer hover:border-primary-500 hover:shadow-md transition focus:outline-none focus:ring-2 focus:ring-primary-200">
        <div class="aspect-[16/10] border-b border-gray-200" style="${style}"></div>
        <div class="px-3 py-2.5 flex items-center gap-2">
          <div class="min-w-0 flex-1">
            <div class="flex items-center gap-1 mb-0.5">
              ${fid}
              <span class="inline-block text-[10px] uppercase tracking-wider bg-gray-100 text-gray-600 px-1.5 py-0.5 rounded">${escHtml(d.type||'prototype')}</span>
            </div>
            <div class="text-sm font-medium text-gray-900 truncate">${escHtml(d.name||'Untitled')}</div>
          </div>
          <div class="text-xs text-gray-400 shrink-0">${fmtDate(d.updated_at||d.created_at||0)}</div>
        </div>
      </button>`;
  }).join("");
  g.querySelectorAll(".card").forEach(c => c.addEventListener("click", () => { window.location.href = "/canvas/" + c.dataset.id; }));
}

async function refreshDs(){
  const r = await fetch("/api/design-systems"); const j = await r.json();
  $("activeDs").textContent = j.active || "(none)";
  const c = $("dsList");
  if (!j.systems.length){
    c.innerHTML = '<div class="text-xs text-gray-400">No design systems yet.</div>';
    return;
  }
  c.innerHTML = j.systems.map(s => {
    const active = s.name === j.active;
    return `<button type="button" data-name="${escHtml(s.name)}" class="ds-item w-full flex items-center gap-2 px-3 py-2 text-sm rounded-lg ${active ? 'bg-gray-100 border border-gray-200' : 'text-gray-700 hover:bg-gray-50'}">
      <span class="w-3 h-3 rounded border border-gray-200 shrink-0" style="background:${escHtml(s.primary_color||'#cccccc')}"></span>
      <span class="truncate flex-1 text-left">${escHtml(s.name)}</span>
      ${active ? '<span class="text-[10px] uppercase tracking-wider text-primary-600 font-semibold">active</span>' : ''}
    </button>`;
  }).join("");
  c.querySelectorAll(".ds-item").forEach(el => el.addEventListener("click", () => openDsModal(el.dataset.name)));
}

let dsEditing = null;
function openDsModal(name){
  dsEditing = name || null;
  $("dsModalTitle").textContent = name ? `Edit: ${name}` : "New Design System";
  $("dsDelete").style.display = name ? "" : "none";
  $("ds_name").value = ""; $("ds_primary").value = "#cc0000"; $("ds_primary_t").value = "#CC0000";
  $("ds_secondary").value = "#111827"; $("ds_secondary_t").value = "#111827";
  $("ds_font").value = ""; $("ds_company").value = ""; $("ds_logo").value = "";
  if (name){
    fetch("/api/design-systems").then(r=>r.json()).then(j => {
      const s = j.systems.find(x => x.name === name);
      if (!s) return;
      $("ds_name").value = s.name;
      if (s.primary_color){ $("ds_primary").value = s.primary_color; $("ds_primary_t").value = s.primary_color; }
      if (s.secondary_color){ $("ds_secondary").value = s.secondary_color; $("ds_secondary_t").value = s.secondary_color; }
      $("ds_font").value = s.font || "";
      $("ds_company").value = s.company_name || "";
      $("ds_logo").value = s.logo_url || "";
    });
  }
  showModal($("dsModal"));
}
function closeDsModal(){ hideModal($("dsModal")); dsEditing = null; }
$("dsNew").addEventListener("click", () => openDsModal(null));
$("dsCancel").addEventListener("click", closeDsModal);
$("dsX").addEventListener("click", closeDsModal);
$("ds_primary").addEventListener("input", () => $("ds_primary_t").value = $("ds_primary").value);
$("ds_primary_t").addEventListener("change", () => { if(/^#[0-9a-fA-F]{3,8}$/.test($("ds_primary_t").value)) $("ds_primary").value = $("ds_primary_t").value; });
$("ds_secondary").addEventListener("input", () => $("ds_secondary_t").value = $("ds_secondary").value);
$("ds_secondary_t").addEventListener("change", () => { if(/^#[0-9a-fA-F]{3,8}$/.test($("ds_secondary_t").value)) $("ds_secondary").value = $("ds_secondary_t").value; });
$("dsSave").addEventListener("click", async () => {
  const body = { name: $("ds_name").value.trim(), primary_color: $("ds_primary_t").value, secondary_color: $("ds_secondary_t").value, font: $("ds_font").value, company_name: $("ds_company").value, logo_url: $("ds_logo").value };
  if (!body.name){ toast("Name required", "err"); return; }
  const r = await fetch("/api/design-systems", { method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify(body) });
  if (r.ok){ await fetch(`/api/design-systems/${encodeURIComponent(body.name)}/activate`, {method:"POST"}); closeDsModal(); refreshDs(); toast("Saved", "ok"); }
  else toast("Save failed", "err");
});
$("dsDelete").addEventListener("click", async () => {
  if (!dsEditing) return;
  if (!confirm(`Delete design system "${dsEditing}"?`)) return;
  await fetch(`/api/design-systems/${encodeURIComponent(dsEditing)}`, {method:"DELETE"});
  closeDsModal(); refreshDs(); toast("Deleted", "ok");
});

$("genBtn").addEventListener("click", async () => {
  const prompt = $("prompt").value.trim();
  if (!prompt){ toast("Enter a prompt first", "err"); return; }
  const name = $("name").value.trim() || prompt.slice(0, 40);
  showModal($("overlay"));
  $("stageTxt").textContent = "Thinking...";
  $("countTxt").textContent = "0 tokens";
  try {
    const res = await fetch("/api/generate", { method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify({ prompt, name, type: selectedType, fidelity: selectedFidelity }) });
    if (!res.ok){ const t = await res.text(); hideModal($("overlay")); toast("Error: " + t.slice(0,80), "err"); return; }
    const reader = res.body.getReader(); const decoder = new TextDecoder();
    let buf = "", id = null, lastTokens = 0;
    while (true){
      const { done, value } = await reader.read(); if (done) break;
      buf += decoder.decode(value, {stream:true});
      let idx;
      while ((idx = buf.indexOf("\n\n")) !== -1){
        const frame = buf.slice(0, idx); buf = buf.slice(idx+2);
        const line = frame.startsWith("data: ") ? frame.slice(6) : frame;
        if (!line) continue;
        let ev; try { ev = JSON.parse(line); } catch(_){ continue; }
        if (ev.stage === "start"){ id = ev.id; }
        else if (ev.stage === "thinking"){ $("stageTxt").textContent = "Thinking..."; }
        else if (ev.stage === "writing"){ lastTokens = ev.tokens; $("stageTxt").textContent = "Writing HTML..."; $("countTxt").textContent = `${nf.format(lastTokens)} tokens`; }
        else if (ev.stage === "done"){ id = ev.id || id; $("stageTxt").textContent = "Done — opening canvas..."; }
        else if (ev.stage === "error"){ hideModal($("overlay")); toast("Error: " + ev.message, "err"); return; }
      }
    }
    if (id) window.location.href = "/canvas/" + id;
    else { hideModal($("overlay")); toast("No design id returned", "err"); }
  } catch(e){ hideModal($("overlay")); toast("Network error: " + e.message, "err"); }
});

/* greeting */
const hour = new Date().getHours();
$("greet").textContent = hour < 5 ? "Working late." : hour < 12 ? "Good morning." : hour < 17 ? "Good afternoon." : "Good evening.";

paintTabs(); paintFid();
refreshRecent(); refreshDs();
</script>
<script src="https://cdn.jsdelivr.net/npm/flowbite@3/dist/flowbite.min.js"></script>
</body></html>
"""


# ---------------------------------------------------------------------------
# Canvas page
# ---------------------------------------------------------------------------

CANVAS_HTML = r"""<!doctype html>
<html lang="en" class="h-full">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>Nexus Design — Canvas</title>
__SHELL_HEAD__
<style>
  .mode-btn { transition: background .15s, color .15s; }
</style>
</head>
<body class="h-full bg-gray-100 text-gray-900 antialiased">
<div class="flex flex-col h-full">

  <!-- Top nav -->
  <nav class="h-12 bg-white border-b border-gray-200 px-3 flex items-center gap-3 shrink-0">
    <a href="/" class="text-gray-500 hover:text-gray-900 text-sm inline-flex items-center gap-1 px-2 py-1 rounded hover:bg-gray-50">
      <svg class="w-4 h-4" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><path d="M15 18l-6-6 6-6"/></svg>
      Home
    </a>
    <div class="flex items-center gap-2 min-w-0 max-w-[40%]">
      <span class="font-semibold text-gray-900 truncate" id="projName">…</span>
      <span id="projType" class="bg-gray-100 text-gray-600 text-[10px] font-medium uppercase tracking-wider px-2 py-0.5 rounded"></span>
    </div>
    <div class="flex-1 flex items-center justify-center">
      <div class="inline-flex rounded-lg border border-gray-200 overflow-hidden shadow-sm" role="group">
        <button id="modeEdit"    class="mode-btn px-3 py-1.5 text-sm font-medium inline-flex items-center gap-1.5">✎ Edit</button>
        <button id="modeComment" class="mode-btn px-3 py-1.5 text-sm font-medium inline-flex items-center gap-1.5 border-l border-gray-200">💬 Comment</button>
        <button id="modeDraw"    class="mode-btn px-3 py-1.5 text-sm font-medium inline-flex items-center gap-1.5 border-l border-gray-200">✏ Draw</button>
      </div>
    </div>
    <div class="relative">
      <button id="expOpen" data-dropdown-toggle="expMenu" data-dropdown-placement="bottom-end" class="text-white bg-primary-600 hover:bg-primary-700 focus:ring-4 focus:outline-none focus:ring-primary-200 font-medium rounded-lg text-sm px-4 py-1.5 inline-flex items-center gap-1">
        Export
        <svg class="w-3 h-3" fill="currentColor" viewBox="0 0 20 20"><path d="M5.23 7.21a.75.75 0 011.06.02L10 11.084l3.71-3.854a.75.75 0 111.08 1.04l-4.24 4.4a.75.75 0 01-1.08 0l-4.24-4.4a.75.75 0 01.02-1.06z"/></svg>
      </button>
      <div id="expMenu" class="z-50 hidden bg-white divide-y divide-gray-100 rounded-lg shadow-lg w-52 border border-gray-200">
        <ul class="py-1 text-sm text-gray-700">
          <li><button data-f="html" class="block w-full text-left px-4 py-2 hover:bg-gray-50">HTML</button></li>
          <li><button data-f="pdf"  class="block w-full text-left px-4 py-2 hover:bg-gray-50">PDF</button></li>
          <li><button data-f="pptx" class="block w-full text-left px-4 py-2 hover:bg-gray-50">PPTX <span class="text-[10px] text-gray-400 uppercase ml-1">slides only</span></button></li>
        </ul>
      </div>
    </div>
  </nav>

  <!-- Main -->
  <div class="flex-1 grid min-h-0" style="grid-template-columns:320px 1fr 320px;">
    <!-- Left: chat + refine input -->
    <aside class="bg-white border-r border-gray-200 flex flex-col min-w-0">
      <div id="chat" class="flex-1 overflow-y-auto p-4 space-y-2"></div>
      <div class="border-t border-gray-200 p-3 space-y-2">
        <textarea id="prompt" rows="3" placeholder="Refine the design — 'make the header bigger', 'swap the accent color', ..." class="block w-full text-sm text-gray-900 bg-white rounded-lg border border-gray-200 focus:ring-primary-500 focus:border-primary-500 p-3 resize-none"></textarea>
        <div class="flex gap-2">
          <button id="refBtn" class="flex-1 text-white bg-primary-600 hover:bg-primary-700 focus:ring-4 focus:outline-none focus:ring-primary-200 font-medium rounded-lg text-sm px-4 py-2 inline-flex items-center justify-center gap-1.5">
            Refine
            <svg class="w-3.5 h-3.5" fill="none" stroke="currentColor" stroke-width="2" viewBox="0 0 24 24"><path d="M5 12h14M13 5l7 7-7 7"/></svg>
          </button>
          <button id="applyCmts" class="text-gray-700 bg-white border border-gray-200 hover:bg-gray-50 font-medium rounded-lg text-sm px-3 py-2">Apply Comments</button>
        </div>
      </div>
    </aside>

    <!-- Center preview -->
    <main class="bg-gray-100 relative overflow-hidden">
      <div class="absolute inset-4 rounded-xl overflow-hidden border border-gray-200 bg-white shadow-sm">
        <iframe id="frame" class="absolute inset-0 w-full h-full border-0 opacity-0 transition-opacity duration-300" sandbox="allow-scripts allow-same-origin"></iframe>
        <svg id="drawLayer" xmlns="http://www.w3.org/2000/svg" class="absolute inset-0 pointer-events-none"></svg>
        <div id="overlay" class="absolute inset-0 bg-white/95 flex flex-col items-center justify-center gap-4 transition-opacity duration-300">
          <div id="bigspin" class="w-12 h-12 border-4 border-gray-200 border-t-primary-600 rounded-full animate-spin"></div>
          <div id="overlayText" class="text-gray-600 text-sm">Loading…</div>
        </div>
      </div>
    </main>

    <!-- Right: tweaks / comments / inspector -->
    <aside class="bg-white border-l border-gray-200 overflow-y-auto divide-y divide-gray-200 min-w-0">
      <div class="p-4 space-y-3">
        <h3 class="text-xs font-semibold text-gray-500 uppercase tracking-wider flex items-center gap-2">Tweaks <span id="tweaksSub" class="font-normal text-gray-400 normal-case tracking-normal"></span></h3>
        <div id="tweaks" class="space-y-3"></div>
      </div>
      <div class="p-4 space-y-3 hidden" id="commentsSec">
        <h3 class="text-xs font-semibold text-gray-500 uppercase tracking-wider">Comments</h3>
        <div id="comments" class="space-y-2"></div>
        <button id="clearCmts" class="text-gray-700 bg-white border border-gray-200 hover:bg-gray-50 font-medium rounded-lg text-sm px-3 py-1.5">Clear pins</button>
      </div>
      <div class="p-4 space-y-2 flex-1">
        <h3 class="text-xs font-semibold text-gray-500 uppercase tracking-wider">Inspector</h3>
        <div id="inspector" class="text-xs text-gray-500">Click an element in the preview to inspect.</div>
      </div>
    </aside>
  </div>

  <!-- Status bar -->
  <div class="h-8 bg-white border-t border-gray-200 flex items-center px-4 text-xs text-gray-500 gap-2 shrink-0" id="statusBar">
    <span id="statusDot" class="w-2 h-2 rounded-full bg-gray-300"></span>
    <span id="statusText">Idle</span>
    <span id="statusCount" class="ml-auto font-mono text-gray-400"></span>
  </div>
</div>

<!-- Context menu -->
<div id="ctx" class="fixed hidden bg-white border border-gray-200 rounded-lg shadow-xl py-1 z-[60] min-w-[190px]">
  <button data-cmd="duplicate" class="block w-full text-left px-3 py-1.5 text-sm text-gray-700 hover:bg-gray-50">Duplicate</button>
  <button data-cmd="delete"    class="block w-full text-left px-3 py-1.5 text-sm text-red-600 hover:bg-red-50">Delete</button>
  <hr class="my-1 border-gray-200">
  <label class="flex items-center justify-between px-3 py-1.5 text-sm text-gray-700 hover:bg-gray-50 cursor-pointer"><span>Text color</span><input type="color" id="ctxColorText" class="w-6 h-5 border-0 p-0 cursor-pointer bg-transparent"></label>
  <label class="flex items-center justify-between px-3 py-1.5 text-sm text-gray-700 hover:bg-gray-50 cursor-pointer"><span>Background</span><input type="color" id="ctxColorBg" class="w-6 h-5 border-0 p-0 cursor-pointer bg-transparent"></label>
</div>

<!-- Toast -->
<div id="toast" class="fixed bottom-5 right-5 bg-gray-900 text-white text-sm px-4 py-2.5 rounded-lg shadow-xl opacity-0 translate-y-2 transition-all duration-200 pointer-events-none z-[60] font-medium"></div>

<script>
const $ = (id) => document.getElementById(id);
const nf = new Intl.NumberFormat();
const ID = window.location.pathname.split("/").pop();
let currentHtml = "";
let generating = false;
let commentMode = false, drawMode = false;
let comments = [];
let tweaks = [];
let pendingCtxTarget = null;

function toast(msg, kind){
  const t = $("toast"); t.textContent = msg;
  t.classList.remove("bg-gray-900","bg-red-600","bg-green-600","opacity-0","translate-y-2","pointer-events-none");
  t.classList.add(kind==="err" ? "bg-red-600" : kind==="ok" ? "bg-green-600" : "bg-gray-900", "opacity-100","translate-y-0");
  clearTimeout(toast._t);
  toast._t = setTimeout(() => { t.classList.remove("opacity-100","translate-y-0"); t.classList.add("opacity-0","translate-y-2","pointer-events-none"); }, 2400);
}
function escHtml(s){ return String(s).replace(/[&<>\"']/g, c => ({"&":"&amp;","<":"&lt;",">":"&gt;","\"":"&quot;","'":"&#39;"}[c])); }

function setDot(mode){
  const d = $("statusDot");
  d.className = "w-2 h-2 rounded-full";
  if (mode === "active") d.classList.add("bg-primary-600","animate-pulse");
  else if (mode === "done") d.classList.add("bg-green-500");
  else if (mode === "err") d.classList.add("bg-red-500");
  else d.classList.add("bg-gray-300");
}
function setBar(text, mode, count){ $("statusText").textContent = text; setDot(mode || ""); $("statusCount").textContent = count ? `${count} tokens` : ""; }
function setOverlay(show, txt){
  const o = $("overlay");
  if (show) o.classList.remove("opacity-0","pointer-events-none"); else o.classList.add("opacity-0","pointer-events-none");
  if (txt) $("overlayText").textContent = txt;
  $("bigspin").style.display = show ? "block" : "none";
}

/* mode buttons styling */
function paintModes(){
  const buttons = [["modeEdit", !commentMode && !drawMode], ["modeComment", commentMode], ["modeDraw", drawMode]];
  for (const [id, on] of buttons){
    const b = $(id);
    b.classList.remove("bg-primary-600","text-white","bg-white","text-gray-700","hover:bg-gray-50");
    b.classList.add(on ? "bg-primary-600" : "bg-white", on ? "text-white" : "text-gray-700", on ? "" : "hover:bg-gray-50");
  }
  $("drawLayer").classList.toggle("pointer-events-auto", drawMode);
  $("drawLayer").style.cursor = drawMode ? "crosshair" : "";
}

/* ---------- injected iframe studio script ---------- */
const INJECTED = `
<script id="__nexus-studio">
(function(){
  const SEL_TEXT = "h1,h2,h3,h4,h5,h6,p,button,span,li,a,strong,em,blockquote,figcaption,label,td,th";
  let mode = { edit: true, comment: false };
  let pinCounter = 0;
  let dragging = null;
  function ensureTweakStyle(){ let s = document.getElementById("__nexus-tweaks-style"); if (!s){ s = document.createElement("style"); s.id="__nexus-tweaks-style"; document.head.appendChild(s);} return s; }
  function cssPath(el){
    if (!(el instanceof Element)) return "";
    const parts = [];
    while (el && el.nodeType === 1 && parts.length < 6){
      let part = el.nodeName.toLowerCase();
      if (el.id){ part += "#" + el.id; parts.unshift(part); break; }
      const cls = (el.className && typeof el.className === "string" ? el.className.trim().split(/\\s+/).slice(0,2).join(".") : "");
      if (cls) part += "." + cls;
      const parent = el.parentNode;
      if (parent){ const sib = Array.from(parent.children).filter(c => c.nodeName === el.nodeName); if (sib.length > 1) part += ":nth-of-type(" + (sib.indexOf(el)+1) + ")"; }
      parts.unshift(part); el = el.parentElement;
    }
    return parts.join(" > ");
  }
  function notifyHtml(){
    const clean = ("<!doctype html>\\n" + document.documentElement.outerHTML)
      .replace(/<script id=\"__nexus-studio\"[\\s\\S]*?<\\/script>/i, "")
      .replace(/<style id=\"__nexus-tweaks-style\"[\\s\\S]*?<\\/style>/i, "")
      .replace(/<div class=\"__nexus-pin\"[\\s\\S]*?<\\/div>/gi, "");
    parent.postMessage({type:"nexus:html", html: clean}, "*");
  }
  function wireEl(el){
    if (el.__nexus) return; el.__nexus = true;
    el.addEventListener("dblclick", (e) => {
      if (!mode.edit || mode.comment) return;
      if (!el.matches(SEL_TEXT)) return;
      e.preventDefault(); e.stopPropagation();
      el.setAttribute("contenteditable", "true");
      el.style.outline = "2px solid #cc0000"; el.style.outlineOffset = "2px";
      el.focus();
    });
    el.addEventListener("blur", () => {
      if (el.getAttribute("contenteditable") === "true"){
        el.removeAttribute("contenteditable"); el.style.outline=""; el.style.outlineOffset="";
        notifyHtml();
      }
    }, true);
    el.addEventListener("keydown", (e) => {
      if (el.getAttribute("contenteditable") === "true"){
        if (e.key === "Enter" && !e.shiftKey && el.tagName !== "TEXTAREA"){ e.preventDefault(); el.blur(); }
        if (e.key === "Escape"){ el.blur(); }
      }
    });
    el.addEventListener("click", (e) => {
      if (mode.comment) return;
      if (el.getAttribute("contenteditable") === "true") return;
      parent.postMessage({type:"nexus:inspect", tag: el.tagName.toLowerCase(), cls: el.className||"", path: cssPath(el)}, "*");
    });
    el.addEventListener("mousedown", (e) => {
      if (!mode.edit || mode.comment) return;
      if (!e.altKey) return;
      if (el === document.body || el === document.documentElement) return;
      e.preventDefault(); e.stopPropagation();
      const rect = el.getBoundingClientRect();
      const cs = getComputedStyle(el);
      if (cs.position === "static" || cs.position === ""){
        el.style.position = "absolute"; el.style.left = (el.offsetLeft) + "px"; el.style.top = (el.offsetTop) + "px";
      }
      const startX = e.clientX, startY = e.clientY;
      const startLeft = parseFloat(el.style.left) || 0, startTop = parseFloat(el.style.top) || 0;
      dragging = el;
      function move(ev){ el.style.left = (startLeft + ev.clientX - startX) + "px"; el.style.top = (startTop + ev.clientY - startY) + "px"; }
      function up(){ window.removeEventListener("mousemove", move); window.removeEventListener("mouseup", up); dragging = null; notifyHtml(); }
      window.addEventListener("mousemove", move); window.addEventListener("mouseup", up);
    });
    el.addEventListener("contextmenu", (e) => {
      if (mode.comment) return;
      e.preventDefault(); e.stopPropagation();
      parent.postMessage({type:"nexus:ctxmenu", x: e.clientX, y: e.clientY, path: cssPath(el), tag: el.tagName.toLowerCase()}, "*");
    });
  }
  function wireAll(){ document.querySelectorAll("*").forEach(wireEl); }
  const mo = new MutationObserver(wireAll);
  mo.observe(document.documentElement, { childList: true, subtree: true });
  wireAll();
  document.addEventListener("click", (e) => {
    if (!mode.comment) return;
    e.preventDefault(); e.stopPropagation();
    pinCounter += 1;
    const n = pinCounter;
    const pin = document.createElement("div");
    pin.className = "__nexus-pin"; pin.textContent = n;
    pin.style.cssText = "position:absolute; left:"+(e.pageX-14)+"px; top:"+(e.pageY-14)+"px; width:28px; height:28px; border-radius:50%; background:#cc0000; color:white; display:flex; align-items:center; justify-content:center; font:700 13px/1 system-ui; z-index:999999; box-shadow:0 4px 14px rgba(0,0,0,.45); cursor:default; user-select:none;";
    document.body.appendChild(pin);
    parent.postMessage({type:"nexus:comment", n}, "*");
  }, true);
  window.addEventListener("message", (ev) => {
    const m = ev.data||{};
    if (m.type === "nexus:mode"){ mode.edit = !!m.edit; mode.comment = !!m.comment; document.body.style.cursor = m.comment ? "crosshair" : ""; }
    else if (m.type === "nexus:clear-pins"){ pinCounter = 0; document.querySelectorAll(".__nexus-pin").forEach(el => el.remove()); }
    else if (m.type === "nexus:tweak-css"){ ensureTweakStyle().textContent = m.css||""; }
    else if (m.type === "nexus:request-html"){ notifyHtml(); }
    else if (m.type === "nexus:ctx-action"){
      try {
        const el = document.querySelector(m.path); if (!el) return;
        if (m.cmd === "delete"){ el.remove(); notifyHtml(); }
        else if (m.cmd === "duplicate"){ const c = el.cloneNode(true); el.parentNode.insertBefore(c, el.nextSibling); notifyHtml(); }
        else if (m.cmd === "color-text"){ el.style.color = m.value; notifyHtml(); }
        else if (m.cmd === "color-bg"){ el.style.backgroundColor = m.value; notifyHtml(); }
      } catch(_){}
    }
  });
  parent.postMessage({type:"nexus:ready"}, "*");
})();
<\/script>`;

function injectStudio(html){
  const lower = html.toLowerCase();
  const body = lower.lastIndexOf("</body>");
  if (body !== -1) return html.slice(0, body) + INJECTED + html.slice(body);
  const close = lower.lastIndexOf("</html>");
  if (close !== -1) return html.slice(0, close) + INJECTED + html.slice(close);
  return html + INJECTED;
}

function showHtml(html){
  currentHtml = html;
  const f = $("frame"); f.classList.remove("opacity-100"); f.classList.add("opacity-0");
  const injected = injectStudio(html);
  const onLoad = () => {
    f.removeEventListener("load", onLoad);
    requestAnimationFrame(() => {
      f.classList.remove("opacity-0");
      f.classList.add("opacity-100");
      setOverlay(false);
      setBar("Ready", "done", "");
      syncModeIntoIframe();
      applyTweakCss();
    });
  };
  f.addEventListener("load", onLoad);
  f.srcdoc = injected;
}

function syncModeIntoIframe(){
  const f = $("frame"); if (!f.contentWindow) return;
  f.contentWindow.postMessage({type:"nexus:mode", edit: true, comment: commentMode}, "*");
}

async function load(){
  const r = await fetch(`/api/designs/${ID}`);
  if (!r.ok){ toast("Design not found", "err"); return; }
  const j = await r.json();
  $("projName").textContent = j.name || "Untitled";
  $("projType").textContent = j.type || "prototype";
  document.title = `${j.name || "Nexus"} — Nexus Design`;
  renderHistory(j.history || []);
  if (j.html){ showHtml(j.html); await loadTweaks(j.html); }
  else { setBar("No HTML yet", "", ""); setOverlay(true, "Waiting for design..."); }
}

function renderHistory(hist){
  $("chat").innerHTML = hist.map(h => `
    <div class="bg-gray-50 border border-gray-200 rounded-lg p-3 text-sm leading-relaxed">
      <span class="inline-block text-[10px] uppercase tracking-wider text-gray-500 mr-2 font-semibold">${escHtml(h.kind)}</span>
      <span class="text-gray-800">${escHtml(h.prompt)}</span>
    </div>`).join("");
  $("chat").scrollTop = $("chat").scrollHeight;
}

/* tweaks */
function applyTweakCss(){
  const css = tweaks.map(t => (t.css||"").replaceAll("{value}", String(t.value))).join("\n");
  const f = $("frame"); if (f.contentWindow) f.contentWindow.postMessage({type:"nexus:tweak-css", css}, "*");
}
function renderTweaks(){
  const c = $("tweaks");
  if (!tweaks.length){ c.innerHTML = '<div class="text-xs text-gray-400">No tweaks yet.</div>'; return; }
  $("tweaksSub").textContent = `· ${tweaks.length}`;
  c.innerHTML = "";
  tweaks.forEach(t => {
    const row = document.createElement("div"); row.className = "space-y-1.5";
    if (t.type === "color"){
      row.innerHTML = `
        <div class="flex justify-between items-center text-xs">
          <span class="font-medium text-gray-700">${escHtml(t.label||t.id)}</span>
          <span class="val font-mono text-gray-400">${escHtml(String(t.value))}</span>
        </div>
        <div class="flex gap-2 items-center">
          <input type="color" value="${escHtml(String(t.value))}" class="w-9 h-8 border border-gray-200 rounded cursor-pointer p-0.5 bg-white shrink-0">
          <input type="text" value="${escHtml(String(t.value))}" class="flex-1 min-w-0 bg-white border border-gray-200 text-gray-900 text-xs rounded-lg focus:ring-primary-500 focus:border-primary-500 px-2 py-1.5">
        </div>`;
      const [pk, tx] = row.querySelectorAll("input");
      pk.addEventListener("input", () => { t.value = pk.value; tx.value = pk.value; row.querySelector(".val").textContent = t.value; applyTweakCss(); });
      tx.addEventListener("change", () => { t.value = tx.value; pk.value = tx.value; row.querySelector(".val").textContent = t.value; applyTweakCss(); });
    } else {
      const unit = t.unit||"";
      row.innerHTML = `
        <div class="flex justify-between items-center text-xs">
          <span class="font-medium text-gray-700">${escHtml(t.label||t.id)}</span>
          <span class="val font-mono text-gray-400">${t.value}${escHtml(unit)}</span>
        </div>
        <input type="range" min="${Number(t.min??0)}" max="${Number(t.max??64)}" step="${Number(t.step??1)}" value="${Number(t.value)}" class="w-full accent-primary-600">`;
      const r = row.querySelector("input"); r.addEventListener("input", () => { t.value = Number(r.value); row.querySelector(".val").textContent = `${t.value}${unit}`; applyTweakCss(); });
    }
    c.appendChild(row);
  });
  applyTweakCss();
}
async function loadTweaks(html){
  $("tweaksSub").textContent = "· generating...";
  $("tweaks").innerHTML = '<div class="text-xs text-gray-400">Analyzing design...</div>';
  try {
    const r = await fetch("/api/tweaks", { method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify({ html }) });
    if (!r.ok){ $("tweaksSub").textContent = "· failed"; $("tweaks").innerHTML = '<div class="text-xs text-red-500">Tweaks unavailable.</div>'; return; }
    const j = await r.json();
    tweaks = (j.tweaks||[]).filter(t => t && t.id && (t.type==="color"||t.type==="range") && t.css);
    renderTweaks();
  } catch(e){ $("tweaksSub").textContent = "· error"; }
}

/* comments */
function renderComments(){
  const sec = $("commentsSec");
  if (!comments.length){ sec.classList.add("hidden"); $("comments").innerHTML = ""; return; }
  sec.classList.remove("hidden");
  $("comments").innerHTML = "";
  comments.forEach(c => {
    const row = document.createElement("div");
    row.className = "flex items-start gap-2 bg-gray-50 border border-gray-200 rounded-lg p-2";
    row.innerHTML = `
      <div class="w-6 h-6 rounded-full bg-primary-600 text-white text-xs font-semibold flex items-center justify-center shrink-0">${c.n}</div>
      <input type="text" placeholder="Note for pin ${c.n}..." value="${escHtml(c.text||"")}" class="flex-1 bg-transparent border-0 outline-none text-sm text-gray-900 px-0 py-0.5 focus:ring-0">
      <button class="x text-gray-400 hover:text-gray-900 px-1 text-sm" title="Remove">×</button>`;
    const inp = row.querySelector("input"); inp.addEventListener("input", () => c.text = inp.value);
    row.querySelector(".x").addEventListener("click", () => { comments = comments.filter(x => x.n !== c.n); renderComments(); });
    $("comments").appendChild(row);
  });
}
function clearComments(){ comments = []; renderComments(); const f = $("frame"); if (f.contentWindow) f.contentWindow.postMessage({type:"nexus:clear-pins"}, "*"); }

$("modeComment").addEventListener("click", () => { commentMode = !commentMode; drawMode = false; paintModes(); syncModeIntoIframe(); });
$("modeEdit").addEventListener("click", () => { commentMode = false; drawMode = false; paintModes(); syncModeIntoIframe(); });
$("modeDraw").addEventListener("click", () => { drawMode = !drawMode; commentMode = false; paintModes(); syncModeIntoIframe(); });
$("clearCmts").addEventListener("click", clearComments);
$("applyCmts").addEventListener("click", () => {
  const notes = comments.filter(c => (c.text||"").trim());
  if (!notes.length){ toast("No comment text", "err"); return; }
  streamRefine(`/api/designs/${ID}/apply-comments`, { comments: notes });
  clearComments();
});

/* refine streaming */
async function streamRefine(path, body){
  if (generating) return;
  generating = true;
  $("refBtn").disabled = true; $("applyCmts").disabled = true;
  $("refBtn").classList.add("opacity-60","cursor-not-allowed"); $("applyCmts").classList.add("opacity-60","cursor-not-allowed");
  setOverlay(true, "Thinking...");
  $("frame").classList.remove("opacity-100"); $("frame").classList.add("opacity-0");
  setBar("Thinking...", "active", "");
  let res;
  try { res = await fetch(path, { method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify(body) }); }
  catch(e){ setBar("Error", "err", ""); setOverlay(false); generating = false; $("refBtn").disabled = false; $("applyCmts").disabled = false; $("refBtn").classList.remove("opacity-60","cursor-not-allowed"); $("applyCmts").classList.remove("opacity-60","cursor-not-allowed"); toast(e.message, "err"); return; }
  if (!res.ok){ const t = await res.text(); setBar("Error", "err", ""); setOverlay(false); generating = false; $("refBtn").disabled = false; $("applyCmts").disabled = false; $("refBtn").classList.remove("opacity-60","cursor-not-allowed"); $("applyCmts").classList.remove("opacity-60","cursor-not-allowed"); toast(t.slice(0,120), "err"); return; }
  const reader = res.body.getReader(); const decoder = new TextDecoder();
  let buf = "", lastTokens = 0, finalHtml = null;
  while (true){
    const { done, value } = await reader.read(); if (done) break;
    buf += decoder.decode(value, {stream:true});
    let idx;
    while ((idx = buf.indexOf("\n\n")) !== -1){
      const frame = buf.slice(0, idx); buf = buf.slice(idx+2);
      const line = frame.startsWith("data: ") ? frame.slice(6) : frame; if (!line) continue;
      let ev; try { ev = JSON.parse(line); } catch(_){ continue; }
      if (ev.stage === "writing"){ lastTokens = ev.tokens; setBar("Writing HTML...", "active", nf.format(lastTokens)); setOverlay(true, "Writing HTML..."); }
      else if (ev.stage === "done"){ finalHtml = ev.html; setBar("Rendering preview...", "active", nf.format(ev.tokens)); }
      else if (ev.stage === "error"){ setBar("Error", "err", ""); setOverlay(false); generating = false; $("refBtn").disabled = false; $("applyCmts").disabled = false; $("refBtn").classList.remove("opacity-60","cursor-not-allowed"); $("applyCmts").classList.remove("opacity-60","cursor-not-allowed"); toast(ev.message, "err"); return; }
    }
  }
  if (finalHtml != null){ showHtml(finalHtml); loadTweaks(finalHtml); refreshProj(); }
  generating = false; $("refBtn").disabled = false; $("applyCmts").disabled = false; $("refBtn").classList.remove("opacity-60","cursor-not-allowed"); $("applyCmts").classList.remove("opacity-60","cursor-not-allowed");
}
async function refreshProj(){ try { const j = await (await fetch(`/api/designs/${ID}`)).json(); renderHistory(j.history||[]); } catch(_){} }
$("refBtn").addEventListener("click", () => {
  const p = $("prompt").value.trim(); if (!p) return;
  streamRefine(`/api/designs/${ID}/refine`, { prompt: p });
  $("prompt").value = "";
});
$("prompt").addEventListener("keydown", (e) => { if ((e.metaKey || e.ctrlKey) && e.key === "Enter"){ e.preventDefault(); $("refBtn").click(); }});

/* export menu — fall back to manual open if Flowbite not ready */
function toggleExp(on){
  const m = $("expMenu");
  if (on === undefined) m.classList.toggle("hidden"); else m.classList.toggle("hidden", !on);
}
$("expOpen").addEventListener("click", (e) => { e.stopPropagation(); toggleExp(); });
document.addEventListener("click", () => toggleExp(false));
$("expMenu").addEventListener("click", (e) => e.stopPropagation());
$("expMenu").querySelectorAll("button").forEach(b => b.addEventListener("click", () => {
  const fmt = b.dataset.f; toggleExp(false);
  window.location.href = `/api/designs/${ID}/export/${fmt}`;
}));

/* draw overlay */
(function(){
  const svg = $("drawLayer");
  let drawing = false, path = null, pts = [];
  function toSvg(e){ const r = svg.getBoundingClientRect(); return {x: e.clientX - r.left, y: e.clientY - r.top}; }
  svg.addEventListener("mousedown", (e) => {
    if (!drawMode) return;
    drawing = true; pts = [toSvg(e)];
    path = document.createElementNS("http://www.w3.org/2000/svg","path");
    path.setAttribute("stroke","#cc0000"); path.setAttribute("stroke-width","3"); path.setAttribute("stroke-linecap","round"); path.setAttribute("stroke-linejoin","round"); path.setAttribute("fill","none");
    svg.appendChild(path);
  });
  svg.addEventListener("mousemove", (e) => {
    if (!drawing) return;
    pts.push(toSvg(e));
    const d = pts.map((p,i) => (i===0?"M":"L") + p.x + "," + p.y).join(" ");
    path.setAttribute("d", d);
  });
  window.addEventListener("mouseup", () => { drawing = false; });
})();

/* context menu bridge */
const ctx = $("ctx");
function openCtx(x, y, path){ pendingCtxTarget = path; ctx.style.left = x + "px"; ctx.style.top = y + "px"; ctx.classList.remove("hidden"); }
function closeCtx(){ ctx.classList.add("hidden"); pendingCtxTarget = null; }
document.addEventListener("click", () => closeCtx());
ctx.addEventListener("click", (e) => e.stopPropagation());
ctx.querySelectorAll("button").forEach(b => b.addEventListener("click", () => {
  const cmd = b.dataset.cmd; if (!pendingCtxTarget) return;
  const f = $("frame"); if (!f.contentWindow){ closeCtx(); return; }
  f.contentWindow.postMessage({type:"nexus:ctx-action", cmd, path: pendingCtxTarget}, "*");
  closeCtx();
}));
["ctxColorText","ctxColorBg"].forEach(id => {
  $(id).addEventListener("change", () => {
    if (!pendingCtxTarget) return;
    const f = $("frame"); if (!f.contentWindow) return;
    const cmd = id === "ctxColorText" ? "color-text" : "color-bg";
    f.contentWindow.postMessage({type:"nexus:ctx-action", cmd, path: pendingCtxTarget, value: $(id).value}, "*");
    closeCtx();
  });
  $(id).addEventListener("click", (e) => e.stopPropagation());
});

/* message bus from iframe */
window.addEventListener("message", (ev) => {
  const m = ev.data||{};
  if (m.type === "nexus:ready"){ syncModeIntoIframe(); }
  else if (m.type === "nexus:html"){
    currentHtml = m.html;
    fetch(`/api/designs/${ID}/html`, { method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify({ html: currentHtml }) });
  }
  else if (m.type === "nexus:inspect"){
    $("inspector").innerHTML = `
      <div class="flex items-center gap-2 mb-1"><span class="text-primary-600 font-semibold text-sm">&lt;${escHtml(m.tag)}&gt;</span></div>
      <div class="text-xs"><span class="text-gray-500 uppercase tracking-wider mr-1">class:</span><span class="text-gray-800 break-all">${escHtml(m.cls||'—')}</span></div>
      <div class="mt-2 text-[11px] text-gray-400 font-mono break-all">${escHtml(m.path||'')}</div>`;
  }
  else if (m.type === "nexus:comment"){
    comments.push({ n: m.n, text: "" }); renderComments();
  }
  else if (m.type === "nexus:ctxmenu"){
    const r = $("frame").getBoundingClientRect();
    openCtx(r.left + m.x, r.top + m.y, m.path);
  }
});

paintModes();
setBar("Idle", "", "");
load();
</script>
<script src="https://cdn.jsdelivr.net/npm/flowbite@3/dist/flowbite.min.js"></script>
</body></html>
"""


# ---------------------------------------------------------------------------
# Route handlers
# ---------------------------------------------------------------------------

@app.get("/", response_class=HTMLResponse)
async def landing():
    return HTMLResponse(LANDING_HTML.replace("__SHELL_HEAD__", SHELL_HEAD))


@app.get("/canvas/{id}", response_class=HTMLResponse)
async def canvas(id: str):
    _load_metadata(id)
    return HTMLResponse(CANVAS_HTML.replace("__SHELL_HEAD__", SHELL_HEAD))


@app.get("/healthz")
async def healthz():
    return {"ok": True}


def main() -> None:
    import uvicorn
    DESIGNS_DIR.mkdir(parents=True, exist_ok=True)
    _ensure_systems_dir()
    uvicorn.run(app, host=HOST, port=PORT, log_level="info")


if __name__ == "__main__":
    main()
